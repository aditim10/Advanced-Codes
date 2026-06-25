"""Generate AnimeApp.pptx — a showcase presentation for the AnimeApp project.

Focus: what we learned + the outcome (not a code walkthrough).
Slide 2 uses real PowerPoint "zoom/pop" entrance animations that cascade in.

Run with:  .pptx-venv/bin/python build_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml

# ---------------------------------------------------------------- palette
BG_DARK = RGBColor(0x0F, 0x14, 0x2E)
BG_PANEL = RGBColor(0x1A, 0x22, 0x44)
ACCENT = RGBColor(0x7C, 0x5C, 0xFF)
ACCENT2 = RGBColor(0x2E, 0xD5, 0xC4)
TEXT_LIGHT = RGBColor(0xF2, 0xF4, 0xFB)
TEXT_MUTED = RGBColor(0xA9, 0xB2, 0xD0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Helvetica Neue"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def add_bg(slide, color=BG_DARK):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    slide.shapes._spTree.remove(shp._element)
    slide.shapes._spTree.insert(2, shp._element)
    return shp


def add_rect(slide, x, y, w, h, color, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size, color=TEXT_LIGHT, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.name = font
    f.color.rgb = color
    return tb


def header(slide, number, title, kicker=None):
    add_rect(slide, Inches(0.6), Inches(0.55), Inches(0.16), Inches(0.95), ACCENT)
    if kicker:
        add_text(slide, Inches(0.95), Inches(0.5), Inches(11), Inches(0.4),
                 kicker.upper(), 13, ACCENT2, bold=True)
        ty = Inches(0.86)
    else:
        ty = Inches(0.66)
    add_text(slide, Inches(0.95), ty, Inches(11.4), Inches(0.9), title, 30,
             TEXT_LIGHT, bold=True)
    add_text(slide, Inches(12.4), Inches(6.95), Inches(0.7), Inches(0.4),
             str(number), 12, TEXT_MUTED, align=PP_ALIGN.RIGHT)
    add_rect(slide, Inches(0.95), Inches(1.62), Inches(11.45), Pt(2), BG_PANEL)


def bullets(slide, items, x, y, w, h, size=16, gap=8):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        level = 0
        text = item
        if isinstance(item, tuple):
            level, text = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.level = level
        run = p.add_run()
        if level == 0:
            run.text = "▸  " + text
            run.font.size = Pt(size)
            run.font.color.rgb = TEXT_LIGHT
        else:
            run.text = "•  " + text
            run.font.size = Pt(size - 2)
            run.font.color.rgb = TEXT_MUTED
        run.font.name = FONT
    return tb


def outcome(slide, text, y=Inches(6.3)):
    add_rect(slide, Inches(0.95), y, Inches(0.16), Inches(0.5), ACCENT2)
    add_text(slide, Inches(1.25), y - Inches(0.03), Inches(11.2), Inches(0.6),
             "Outcome:  " + text, 15, ACCENT2, bold=True, italic=True,
             anchor=MSO_ANCHOR.MIDDLE)


def card(slide, x, y, w, h, title, lines, title_color=ACCENT2):
    add_rect(slide, x, y, w, h, BG_PANEL)
    add_rect(slide, x, y, w, Inches(0.09), title_color)
    add_text(slide, x + Inches(0.22), y + Inches(0.2), w - Inches(0.4),
             Inches(0.5), title, 16, WHITE, bold=True)
    tb = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.72),
                                  w - Inches(0.44), h - Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(5)
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(12.5)
        r.font.color.rgb = TEXT_MUTED
        r.font.name = FONT


def pill(slide, x, y, w, h, title, sub, accent):
    """A single auto-shape (so it can be animated as one object)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shp.adjustments[0] = 0.10
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = BG_PANEL
    shp.line.color.rgb = accent
    shp.line.width = Pt(2.5)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(19)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = FONT
    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size = Pt(11.5)
        r2.font.color.rgb = ACCENT2
        r2.font.name = FONT
    return shp


# ---------------------------------------------------------------- animation
def _zoom_par(spid, cid, delay, node):
    return f'''<p:par>
      <p:cTn id="{cid}" presetID="23" presetClass="entr" presetSubtype="0" fill="hold" grpId="0" nodeType="{node}">
        <p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>
        <p:childTnLst>
          <p:set>
            <p:cBhvr>
              <p:cTn id="{cid+1}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
              <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
              <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
            </p:cBhvr>
            <p:to><p:strVal val="visible"/></p:to>
          </p:set>
          <p:anim calcmode="lin" valueType="num">
            <p:cBhvr additive="base">
              <p:cTn id="{cid+2}" dur="500"/>
              <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
              <p:attrNameLst><p:attrName>ppt_w</p:attrName></p:attrNameLst>
            </p:cBhvr>
            <p:tavLst>
              <p:tav tm="0"><p:val><p:fltVal val="0.05"/></p:val></p:tav>
              <p:tav tm="100000"><p:val><p:fltVal val="1"/></p:val></p:tav>
            </p:tavLst>
          </p:anim>
          <p:anim calcmode="lin" valueType="num">
            <p:cBhvr additive="base">
              <p:cTn id="{cid+3}" dur="500"/>
              <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
              <p:attrNameLst><p:attrName>ppt_h</p:attrName></p:attrNameLst>
            </p:cBhvr>
            <p:tavLst>
              <p:tav tm="0"><p:val><p:fltVal val="0.05"/></p:val></p:tav>
              <p:tav tm="100000"><p:val><p:fltVal val="1"/></p:val></p:tav>
            </p:tavLst>
          </p:anim>
          <p:animEffect transition="in" filter="fade">
            <p:cBhvr>
              <p:cTn id="{cid+4}" dur="500"/>
              <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
            </p:cBhvr>
          </p:animEffect>
        </p:childTnLst>
      </p:cTn>
    </p:par>'''


def add_cascade_zoom(slide, spids, step=300):
    """Attach a one-click cascade of zoom/pop entrance effects to the shapes."""
    effects = []
    cid = 10
    for i, spid in enumerate(spids):
        node = "clickEffect" if i == 0 else "withEffect"
        delay = 0 if i == 0 else i * step
        effects.append(_zoom_par(spid, cid, delay, node))
        cid += 10
    timing = f'''<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
        xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                <p:par>
                  <p:cTn id="3" fill="hold">
                    <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="4" fill="hold">
                          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                          <p:childTnLst>
                            {''.join(effects)}
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
            <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>'''
    slide._element.append(parse_xml(timing))


# ================================================================ SLIDE 1 — Title
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(0.22), ACCENT)
add_rect(s, 0, SLIDE_H - Inches(0.22), SLIDE_W, Inches(0.22), ACCENT2)
add_text(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(0.5),
         "PROJECT SHOWCASE", 16, ACCENT2, bold=True)
add_text(s, Inches(0.95), Inches(2.75), Inches(11.5), Inches(1.3),
         "AnimeApp", 60, WHITE, bold=True)
add_text(s, Inches(1.0), Inches(4.0), Inches(11.3), Inches(0.9),
         "What We Built & What We Learned", 26, TEXT_LIGHT)
add_text(s, Inches(1.0), Inches(4.75), Inches(11.3), Inches(0.7),
         "A native iOS anime app — from concept to a polished, tested product",
         16, TEXT_MUTED, italic=True)

# ================================================================ SLIDE 2 — Key concepts (ANIMATED)
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_text(s, Inches(0.95), Inches(0.55), Inches(11.4), Inches(0.45),
         "WHAT THIS PROJECT IS ALL ABOUT", 14, ACCENT2, bold=True)
add_text(s, Inches(0.95), Inches(0.95), Inches(11.4), Inches(0.8),
         "Key Concepts We Covered", 32, WHITE, bold=True)

concepts = [
    ("UI / UX", "Adaptive · themed · animated", ACCENT2),
    ("Networking Layer", "Typed · resilient · paginated", ACCENT),
    ("Design Patterns & Architecture", "VIP / Clean Swift", ACCENT2),
    ("Industry-Level SDKs", "Reusable building blocks", ACCENT),
    ("Modular Approach", "API Client · Image Caching ·\nAnalytics Adapter · Player SDK", ACCENT2),
    ("Test Cases", "Business logic — full coverage", ACCENT),
]
xs = [Inches(0.7), Inches(4.815), Inches(8.93)]
ys = [Inches(2.35), Inches(4.7)]
pw = Inches(3.7)
ph = Inches(1.9)
spids = []
for i, (t, sub, col) in enumerate(concepts):
    x = xs[i % 3]
    y = ys[i // 3]
    shp = pill(s, x, y, pw, ph, t, sub, col)
    spids.append(shp.shape_id)
add_cascade_zoom(s, spids, step=320)

# ================================================================ SLIDE 3 — What we built
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, 3, "What We Built", "The Outcome")
add_text(s, Inches(0.95), Inches(1.85), Inches(11.4), Inches(0.6),
         "A complete, API-backed anime browser — shipped end to end.",
         16, TEXT_LIGHT)
cards = [
    ("Login", ["Validated sign-in", "Persistent session"]),
    ("Home Feed", ["Featured banner", "Genre carousels", "Infinite scroll"]),
    ("Search", ["Instant, debounced", "Paginated results"]),
    ("Detail + Trailer", ["Synopsis & cast", "In-app trailer", "Character profiles"]),
    ("Theming", ["Light / Dark / Auto", "Instant switching"]),
]
cx = Inches(0.95)
cw = Inches(2.22)
for i, (t, l) in enumerate(cards):
    card(s, cx + i * (cw + Inches(0.11)), Inches(2.7), cw, Inches(2.6), t, l)
outcome(s, "A real product feel — not just a demo.")

# ================================================================ SLIDE 4 — UI / UX
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, 4, "UI / UX", "Concept 1")
bullets(s, [
    "Adaptive layouts that adjust across device sizes and orientations",
    "Light / Dark / Auto theming with instant switching",
    "Smooth entrance animations and shimmer loading states",
    "Flicker-free list updates using diffable data sources",
    "Localization-ready text and accessibility-minded design",
], Inches(0.95), Inches(2.05), Inches(11.4), Inches(3.8), size=18, gap=16)
outcome(s, "A polished, responsive experience that feels truly native.")

# ================================================================ SLIDE 5 — Networking
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, 5, "Networking Layer", "Concept 2")
bullets(s, [
    "A typed, reusable API client covering every endpoint",
    "Automatic retries and backoff for flaky networks and rate limits",
    "Seamless pagination powering infinite scroll",
    "Clean separation between raw API data and clean app models",
], Inches(0.95), Inches(2.05), Inches(11.4), Inches(3.6), size=18, gap=18)
outcome(s, "Reliable data loading that gracefully handles real-world conditions.")

# ================================================================ SLIDE 6 — Architecture & patterns
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, 6, "Design Patterns & Architecture", "Concept 3")
bullets(s, [
    "VIP / Clean Swift — a clear separation of concerns",
    "Every feature follows the same predictable structure",
    "Proven patterns: Singleton, Facade, Adapter, Observer, Strategy",
    "Protocol-driven boundaries make code easy to test and extend",
], Inches(0.95), Inches(2.05), Inches(6.9), Inches(3.6), size=16, gap=15)
labels = [("View", ACCENT), ("Interactor", ACCENT2), ("Presenter", ACCENT)]
for i, (t, c) in enumerate(labels):
    y = Inches(2.15) + i * Inches(1.15)
    add_rect(s, Inches(8.2), y, Inches(4.3), Inches(0.9), BG_PANEL, line=c)
    add_text(s, Inches(8.2), y, Inches(4.3), Inches(0.9), t, 16, WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(8.2), Inches(5.6), Inches(4.3), Inches(0.4),
         "predictable, testable flow", 12.5, TEXT_MUTED,
         align=PP_ALIGN.CENTER, italic=True)
outcome(s, "A scalable codebase that is easy to grow and maintain.")

# ================================================================ SLIDE 7 — Engineering principles
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, 7, "Built to Last", "Engineering Principles")
add_text(s, Inches(0.95), Inches(1.85), Inches(11.4), Inches(0.5),
         "The whole app is designed around three core principles of good software.",
         15, TEXT_LIGHT)
principles = [
    ("MAINTAINABLE",
     ["Clear separation of concerns",
      "Consistent module structure",
      "No massive view controllers",
      "Easy to read and change"],
     ACCENT2),
    ("SCALABLE",
     ["Add features without breaking others",
      "Reusable, modular SPM packages",
      "Same pattern for every screen",
      "Grows with the product"],
     ACCENT),
    ("RELIABLE",
     ["Resilient networking & retries",
      "Full business-logic test coverage",
      "Predictable data flow",
      "Safe to ship and refactor"],
     ACCENT2),
]
pcw = Inches(3.7)
for i, (t, lines, c) in enumerate(principles):
    x = Inches(0.95) + i * (pcw + Inches(0.18))
    y = Inches(2.55)
    h = Inches(3.5)
    add_rect(s, x, y, pcw, h, BG_PANEL)
    add_rect(s, x, y, pcw, Inches(0.7), c)
    add_text(s, x, y + Inches(0.08), pcw, Inches(0.55), t, 18, BG_DARK, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tb = s.shapes.add_textbox(x + Inches(0.28), y + Inches(0.95),
                              pcw - Inches(0.56), h - Inches(1.1))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(10)
        r = p.add_run()
        r.text = "•  " + ln
        r.font.size = Pt(14)
        r.font.color.rgb = TEXT_LIGHT
        r.font.name = FONT
outcome(s, "Engineering that holds up as the app and team grow.")

# ================================================================ SLIDE 8 — SDKs & modular approach
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, 8, "Industry-Level SDKs & Modular Approach", "Concept 4 & 5")
add_text(s, Inches(0.95), Inches(1.85), Inches(11.4), Inches(0.5),
         "Self-contained, reusable modules — the way real products are built.",
         15, TEXT_LIGHT)
mods = [
    ("API Client", "Typed networking engine", ACCENT2),
    ("Image Downloader & Caching", "Fast, smooth images", ACCENT),
    ("Analytics Adapter", "Pluggable tracking", ACCENT2),
    ("Player SDK", "In-app video & trailers", ACCENT),
]
for i, (t, d, c) in enumerate(mods):
    x = Inches(0.95) + (i % 2) * Inches(5.9)
    y = Inches(2.55) + (i // 2) * Inches(1.55)
    add_rect(s, x, y, Inches(5.5), Inches(1.3), BG_PANEL)
    add_rect(s, x, y, Inches(0.12), Inches(1.3), c)
    add_text(s, x + Inches(0.35), y + Inches(0.2), Inches(5.0), Inches(0.5), t,
             17, WHITE, bold=True)
    add_text(s, x + Inches(0.35), y + Inches(0.72), Inches(5.0), Inches(0.45), d,
             13.5, TEXT_MUTED)
outcome(s, "Reusable building blocks that mirror professional engineering.")

# ================================================================ SLIDE 9 — Analytics adapter
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, 9, "Analytics Adapter", "Modular Approach")
bullets(s, [
    "One event bus feeding many providers (Console, MixPanel, Firebase)",
    "Swap or add analytics tools without touching feature code",
    "Privacy-by-design: anonymous IDs, never personal data",
    "Full funnel tracked: launch → login → browse → search",
], Inches(0.95), Inches(2.05), Inches(11.4), Inches(3.6), size=18, gap=18)
outcome(s, "Insight into user behavior with zero vendor lock-in.")

# ================================================================ SLIDE 10 — Image caching
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, 10, "Image Downloading & Caching", "Modular Approach")
bullets(s, [
    "In-memory caching avoids re-downloading the same images",
    "Smart de-duplication of identical in-flight requests",
    "Off-main-thread decoding for buttery-smooth scrolling",
    "Auto-cancel on cell reuse — never the wrong or stale image",
], Inches(0.95), Inches(2.05), Inches(11.4), Inches(3.6), size=18, gap=18)
outcome(s, "Fast, smooth image-heavy screens.")

# ================================================================ SLIDE 11 — Player SDK
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, 11, "Player SDK", "Modular Approach")
bullets(s, [
    "In-app trailer and video playback",
    "Custom controls, fullscreen, and YouTube trailers",
    "Reusable across any screen in the app",
], Inches(0.95), Inches(2.05), Inches(11.4), Inches(3.0), size=18, gap=20)
outcome(s, "Rich media — without ever leaving the app.")

# ================================================================ SLIDE 12 — Testing
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, 12, "Test Cases — Full Business-Logic Coverage", "Concept 6")
bullets(s, [
    "All business logic is unit-tested (Interactors & Presenters)",
    "Mocks and spies isolate each layer cleanly",
    "Real API-shaped JSON fixtures — tests run with no network",
    "Fast, deterministic, and repeatable",
], Inches(0.95), Inches(2.05), Inches(11.4), Inches(3.6), size=18, gap=18)
outcome(s, "Confidence to ship and refactor safely.")

# ================================================================ SLIDE 13 — Skills gained
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, 13, "Skills & Concepts Gained", "Recap")
skills = [
    "Clean Architecture", "Protocol-Oriented Design", "Swift Concurrency",
    "Modular SPM Packages", "Dependency Injection", "Networking & REST",
    "Caching Strategies", "Analytics & Privacy", "Unit Testing",
    "UI / UX & Theming", "Design Patterns", "Media Playback",
]
cols = 3
cw = Inches(3.85)
ch = Inches(0.95)
for i, c in enumerate(skills):
    col = i % cols
    row = i // cols
    x = Inches(0.95) + col * (cw + Inches(0.18))
    y = Inches(2.0) + row * (ch + Inches(0.18))
    add_rect(s, x, y, cw, ch, BG_PANEL)
    add_rect(s, x, y, Inches(0.1), ch, ACCENT if i % 2 else ACCENT2)
    add_text(s, x + Inches(0.28), y, cw - Inches(0.4), ch, c, 15, TEXT_LIGHT,
             bold=True, anchor=MSO_ANCHOR.MIDDLE)

# ================================================================ SLIDE 14 — Outcome & impact
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, 14, "Outcome & Impact", "Takeaways")
wins = [
    ("Shipped", "A complete, production-style iOS app, end to end"),
    ("Architected", "Modular, testable, and scalable by design"),
    ("Reusable", "SDKs and modules that mirror industry practice"),
    ("Proven", "Strong engineering fundamentals demonstrated"),
]
for i, (t, d) in enumerate(wins):
    y = Inches(2.05) + i * Inches(1.1)
    add_rect(s, Inches(0.95), y, Inches(2.7), Inches(0.85), ACCENT if i % 2 else ACCENT2)
    add_text(s, Inches(0.95), y, Inches(2.7), Inches(0.85), t, 18, WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(3.95), y, Inches(8.6), Inches(0.85), d, 16, TEXT_LIGHT,
             anchor=MSO_ANCHOR.MIDDLE)

# ================================================================ SLIDE 15 — Thank you
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(0.22), ACCENT)
add_rect(s, 0, SLIDE_H - Inches(0.22), SLIDE_W, Inches(0.22), ACCENT2)
add_text(s, Inches(0.95), Inches(2.6), Inches(11.4), Inches(1.2),
         "Thank You", 52, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.95), Inches(4.0), Inches(11.4), Inches(0.7),
         "Questions & Discussion", 24, ACCENT2, bold=True, align=PP_ALIGN.CENTER)

prs.save("AnimeApp.pptx")
print("Saved AnimeApp.pptx with", len(prs.slides._sldIdLst), "slides")
